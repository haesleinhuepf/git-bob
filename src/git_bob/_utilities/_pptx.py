def make_slides(slides_description_json, filename="issue_slides.pptx"):
    """
    Create a PowerPoint slide deck from a JSON description of slides.

    Parameters
    ----------
    slides_description_json : str
        A JSON string representing the slides. Each slide is expected to be a dictionary with
        a 'title' key and a 'content' key, where 'content' is a list of strings.
    filename : str, optional
        The name of the file to save the presentation to. Defaults to 'issue_slides.pptx'.

    """
    import json
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pathlib import Path
    from PIL import Image
    import os
    import re
    from ._config import IMAGE_FILE_ENDINGS

    print("SLIDES_JSON:", slides_description_json)

    # Parse json-encoded slide description
    slides_description_json = re.sub(r'[\x00-\x1f\x7f]', '', slides_description_json)

    slides_data = json.loads(slides_description_json)

    file_location = "slide_template.pptx"
    if not os.path.exists(file_location):
        file_location = Path(__file__).parent / "data" / "slide_template.pptx"

    # Create a presentation
    presentation = Presentation(file_location)

    # determine slide size
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Convert EMU to inches (1 inch = 914400 EMUs)
    width_inch = slide_width / 914400
    height_inch = slide_height / 914400

    top = Inches(2)
    bottom = Inches(1)

    # Iterate through slide data to create slides
    for i, slide_data in enumerate(slides_data):
        # Add a slide with title and content layout
        slide_layout = presentation.slide_layouts[0 if i == 0 else 1] # choose first or second layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.title
        title_box.text = slide_data['title']

        # Calculate width for content columns
        num_columns = len(slide_data['content'])

        # remove all placeholders except the title
        for shape in slide.placeholders:
            if shape != title_box:
                #shape.text = ""
                top = shape.top
                bottom = height_inch - top
                slide.shapes._spTree.remove(shape._element)

        for i, content in enumerate(slide_data['content']):
            content_width = (Inches(width_inch) - top) / num_columns
            content_height = (Inches(height_inch) - top) - bottom
            left = Inches(1) + i * (content_width + Inches(0.1))

            print("Left", left)

            if any([content.endswith(f) for f in IMAGE_FILE_ENDINGS]):
                image_path = content

                # load image using PIL and determine width/height ratio
                image = Image.open(image_path)
                width, height = image.size
                aspect_ratio = width / height

                if content_width / aspect_ratio > content_height - bottom:
                    content_width = (content_height - bottom) / aspect_ratio

                content_box = slide.shapes.add_picture(image_path, left=left, top=top, width=content_width, height=content_width / aspect_ratio)
            else:
                content_box = slide.shapes.add_textbox(left=left, top=top, width=content_width, height=content_height)
                text_frame = content_box.text_frame
                text_frame.text = content
                text_frame.word_wrap = True

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Set the font size for each run in each paragraph
                        run.font.size = Pt(24)  # Set font size to 24 points


    # Save presentation
    presentation.save(filename)