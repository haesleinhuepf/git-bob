# This module provides utility functions for text processing, including functions to remove indentation and outer markdown from text.
import sys
import warnings
from functools import lru_cache
from functools import wraps
from toolz import curry
from ._endpoints import prompt_openai
import os

VISION_SYSTEM_MESSAGE = os.environ.get("VISION_SYSTEM_MESSAGE", "You are a AI-based vision system. You described images professionally and clearly.")

IMAGE_FILE_ENDINGS = [".jpg", ".png", ".gif", ".jpeg"]
TEXT_FILE_ENDINGS = [".txt", ".md", ".csv", ".yml", ".yaml", ".json", ".py", ".java", ".groovy", ".jython", ".md", ".markdown", ".plaintext", ".tex", ".latex", ".txt", ".csv", ".yml", ".yaml", ".json", ".py", ".svg", ".xml"]

POSSBILE_MARKDOWN_FENCES = ["```python", "```Python", "```nextflow", "```java", "```javascript", "```macro", "```groovy", "```cmd"
                           "```jython", "```md", "```markdown", "```plaintext", "```tex", "```latex",
                           "```txt", "```csv", "```yml", "```yaml", "```json", "```JSON", "```py", "```svg", "```xml", "<FILE>", "```"]

def remove_outer_markdown(text):
    """
    Remove outer markdown syntax from the given text.

    Parameters
    ----------
    text : str
        The input text with potential markdown syntax.

    Returns
    -------
    str
        The text with outer markdown syntax removed and stripped.
    """
    text = text.strip("\n").strip(" ")

    possible_beginnings = POSSBILE_MARKDOWN_FENCES

    possible_endings = ["```", "</FILE>"]

    if any([text.startswith(beginning) for beginning in possible_beginnings]) and any([text.endswith(ending) for ending in possible_endings]):

        for beginning in possible_beginnings:
            if text.startswith(beginning):
                text = text[len(beginning):]
                break

        for ending in possible_endings:
            if text.endswith(ending):
                text = text[:-len(ending)]
                break

    text = text.strip("\n")

    return text


class Config:
    llm_name = None
    run_id = None
    repository = None
    issue = None
    running_in_github_ci = None
    running_in_gitlab_ci = None
    git_server_url = "https://github.com/"
    git_utilities = None
    is_pull_request = None
    pull_request = None
    remarks = []


class ErrorReporting:
    status = False


def quick_first_response(repository, issue):
    """
    Response to a comment to the GitHub issue just mentioning that we're on it.

    Parameters
    ----------
    message : str
        The error message to be reported.
    """
    import sys
    import os
    from ._ai_github_utilities import setup_ai_remark


    run_id = os.environ.get("GITHUB_RUN_ID")
    ai_remark = setup_ai_remark()

    # add reaction to issue
    Config.git_utilities.add_reaction_to_last_comment_in_issue(repository, issue, "+1")


def split_content_and_summary(text):
    """
    Split the given text into content and summary.

    Assuming a text consists of a task solution (code, text) and a summary in the last line,
    it splits the text into the content and the summary.

    Parameters
    ----------
    text : str
        The input text containing content and summary.

    Returns
    -------
    tuple
        A tuple containing two elements:
        - str: The content with outer markdown removed.
        - str: The summary.
    """
    text = text.strip("\n").strip()
    temp = text.split("\n")
    summary = temp[-1].strip()
    remaining_content = temp[:-1]
    if len(summary) < 5:
        summary = temp[-2]
        remaining_content = temp[:-2]

    new_content = remove_outer_markdown("\n".join(remaining_content))

    return new_content.strip(), summary.strip()


def erase_outputs_of_code_cells(file_content):
    """
    Erase outputs of code cells in a Jupyter notebook.

    Parameters
    ----------
    notebook : str
        The notebook content as a string.
    """
    import re
    import json

    # removed invalid characters
    clean_file_content = re.sub(r'[\x00-\x1f\x7f]', '', file_content)

    notebook = json.loads(clean_file_content)
    for cell in notebook.get('cells', []):
        if cell.get('cell_type') == 'code':
            cell['outputs'] = []
            cell['execution_count'] = None
        #cell['id'] = None

    notebook["metadata"] = {}

    file_content = json.dumps(notebook, indent=1)
    return file_content


def restore_outputs_of_code_cells(new_content, original_ipynb_file_content):
    """
    Restore outputs of code cells in a Jupyter notebook from another notebook.
    """
    import json
    print("Recovering outputs in ipynb file")
    original_notebook = json.loads(original_ipynb_file_content)
    new_notebook = json.loads(new_content)

    original_code_cells = [cell for cell in original_notebook['cells'] if cell['cell_type'] == 'code']
    new_code_cells = [cell for cell in new_notebook['cells'] if cell['cell_type'] == 'code']

    if len(original_code_cells) != len(new_code_cells):
        raise ValueError("Number of code cells in the original and new notebooks do not match.")
    
    for o_cell, n_cell in zip(original_code_cells, new_code_cells):
        if "\n".join(o_cell['source']).strip() == "\n".join(n_cell['source']).strip():
            print("Original cell content", o_cell)
            print("New cell content", n_cell)
            if "outputs" in o_cell.keys():
                n_cell['outputs'] = o_cell['outputs']
                n_cell['execution_count'] = o_cell['execution_count']
        else:  # if code is different, any future results may be different, too
            raise ValueError("Code cells are different. Cannot restore outputs.")

    new_notebook["metadata"] = original_notebook["metadata"]
    return json.dumps(new_notebook, indent=1)

def text_to_json(text):
    """Converts a string, e.g. a response from an LLM, to a valid JSON object."""
    import json
    import re

    text = re.sub(r'[\x00-\x1f\x7f]', '', text)

    if "[" in text:
        text = "[" +  text.split("[")[1]
    if "]" in text:
        text = text.split("]")[0] + "]"
    text = text.replace("'", "\"")

    print("JSON?:", text)

    return json.loads(text)


def is_github_url(url):
    """
    Check if the given URL is a GitHub URL and determine its type.
    """
    if not str(url).startswith(Config.git_server_url):
        return None
    if "/.github" in url:
        return None
    if ".gitlab-ci.yml" in url:
        return None
    if '/issues/' in url:
        return 'issue'
    elif '/pull/' in url:
        return 'pull_request'
    elif any([url.endswith(f) for f in IMAGE_FILE_ENDINGS]) \
            or url.endswith('.webp') or "user-attachments/assets" in url or url.endswith("?raw=true"):
        return 'image'
    elif url.endswith('.csv') or url.endswith('.xlsx') or url.endswith('.tif') or url.endswith('.zip') or url.endswith('.svg') or url.endswith('.xml'):
        return 'data'
    elif 'blob/' in url:
        return 'file'
    return None


def load_image_from_url(url):
    import urllib
    import io
    from PIL import Image
    from ._logger import Log
    Log().log("Loading image from URL: " + url)

    # Load the bytestream from the URL
    with urllib.request.urlopen(url) as response:
        bytestream = response.read()

    # Open the image from bytestream using PIL
    image = Image.open(io.BytesIO(bytestream))
    return image


def image_to_url(image):
    """
    Convert an image to a URL.
    """
    if isinstance(image, str) and (image.startswith("data:image") or image.startswith("http")):
        return image

    import base64
    import io
    from PIL import Image

    if isinstance(image, str):
        return image

    if isinstance(image, bytes):
        image = Image.open(io.BytesIO(image))

    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    img_str = base64.b64encode(buffered.getvalue()).decode()
    return img_str


def modify_discussion(discussion, prompt_visionlm=prompt_openai):
    import re
    import docx2markdown
    #from ._github_utilities import get_conversation_on_issue, get_diff_of_pull_request, get_file_in_repository

    # Regex to find URLs in the discussion
    url_pattern = r'(https?://[^\s]+)'
    urls = re.findall(url_pattern, discussion)

    # Placeholder for additional content extracted from URLs
    additional_content = {}

    # Process each URL based on its type
    for url in urls:
        if url.endswith(")"): # happens with ![](url) syntax
            url = url[:-1]
        if url.endswith("'"):
            url = url[:-1]
        if url.endswith('"'):
            url = url[:-1]
            
        url_type = is_github_url(url)
        print("URL:", url)
        print("Type:", url_type)
        from datetime import datetime

        current_datetime = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

        if "### File {url} content" in discussion:
            continue

        try:
            if url_type == 'issue':
                parts = url.split('/')
                repo = parts[3] + '/' + parts[4]
                try:
                    issue_number = int(parts[-1])
                except:
                    continue
                additional_content[url] = Config.git_utilities.get_conversation_on_issue(repo, issue_number)
            elif url_type == 'pull_request':
                parts = url.split('/')
                repo = parts[3] + '/' + parts[4]
                try:
                    pr_number = int(parts[-1])
                except:
                    continue

                # Get both the diff and discussion on pull request
                additional_content[url] = (Config.git_utilities.get_conversation_on_issue(repo, pr_number) +
                                           Config.git_utilities.get_diff_of_pull_request(repo, pr_number))
            elif url_type == 'file':
                parts = url.split('/')
                #repo = parts[3] + '/' + parts[4]
                #branch_name = parts[6]
                file_path = '/'.join(parts[7:])
                temp_file_name = current_datetime + "_" + file_path.split('/')[-1]
                url = url.replace("/blob/", "/raw/")
                download_url(url, temp_file_name)
                #file_contents = Config.git_utilities.get_file_in_repository (repo, branch_name, file_path).decoded_content.decode()
                if url.endswith('.ipynb'):
                    file_contents = read_text_file(temp_file_name)
                    file_contents = erase_outputs_of_code_cells(file_contents)
                elif url.endswith('.docx'):
                    docx2markdown.docx_to_markdown(temp_file_name, temp_file_name + ".md")
                    file_contents = read_text_file(temp_file_name + ".md")
                    # remove the file
                    os.remove(temp_file_name + ".md")
                else:
                    file_contents = read_text_file(temp_file_name)

                os.remove(temp_file_name)

                additional_content[url] = file_contents
            elif url_type == 'image':
                image = load_image_from_url(url)
                image_content = prompt_visionlm(VISION_SYSTEM_MESSAGE + "\n\nDescribe this image.", image=image)
                additional_content[url] = image_content
        except Exception as e:
            print(f"Error while processing URL {url}: {e}")

    # read local files
    temp = discussion.replace("\n", " ")
    for potential_filename in temp.split(" "):
        if len(potential_filename) < 4: # too short to be a filename
            continue
        # check if file exists
        if os.path.exists(potential_filename):
            if potential_filename.endswith('.ipynb'):
                file_contents = read_text_file(potential_filename)
                file_contents = erase_outputs_of_code_cells(file_contents)
            elif potential_filename.endswith('.docx'):
                docx2markdown.docx_to_markdown(potential_filename, potential_filename + ".md")
                file_contents = read_text_file(potential_filename + ".md")
                # remove the file
                os.remove(potential_filename + ".md")
            elif any([potential_filename.endswith(f) for f in TEXT_FILE_ENDINGS]):
                file_contents = read_text_file(potential_filename)
            else:
                continue

            additional_content[potential_filename] = file_contents

    # Modify the existing discussion content
    discussion = discussion.replace("\n#", "\n###")
    discussion = re.sub(r'<sup>.*?</sup>', '', discussion)

    # Append the additional content to the discussion before returning
    temp = []
    for k, v in additional_content.items():
        temp = temp + [f"### File {k} content\n\n```\n{v}\n```\n"]

    return discussion + "\n\n" + "\n".join(temp)


def download_url(url, output_path):
    """
    Downloads a file from the specified URL and saves it to the given output path.

    Args:
        url (str): The URL of the file to download.
        output_path (str): The path where the file will be saved.

    Returns:
        None
    """
    print("Downloading ", url, " to ", output_path)
    import requests
    import os
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()  # Raise an HTTPError for bad responses (4xx and 5xx)

        # Write the content to the output file in binary mode
        with open(output_path, 'wb') as file:
            for chunk in response.iter_content(chunk_size=8192):
                file.write(chunk)

        print(f"File downloaded successfully and saved to {output_path}")

    except requests.exceptions.RequestException as e:
        print(f"Error occurred while downloading the file: {e}")


def execute_notebook(notebook_content, timeout=600, kernel_name='python3'):
    """
    Execute a Jupyter notebook and return whether an error occurred.

    Args:
        notebook_content (str): Content of the notebook file as string in json format.
        timeout (int): Timeout in seconds for each cell (default 600).
        kernel_name (str): The kernel to use for execution (default 'python3').

    Returns:
        str: content of the executed notebook.

    """
    import nbformat
    from nbconvert.preprocessors import ExecutePreprocessor
    import jupyter_client
    import traceback

    # Get the list of available kernels
    kernels = jupyter_client.kernelspec.KernelSpecManager().get_all_specs()

    # Print the names of the kernels
    print("Available Jupyter kernels:")
    for name, spec in kernels.items():
        print(f"- {name}")

    # Load the notebook
    notebook = nbformat.reads(notebook_content, as_version=4)

    # Initialize the processor to execute the notebook
    ep = ExecutePreprocessor(timeout=timeout, kernel_name=kernel_name)

    error = None
    try:
        # Execute the notebook
        ep.preprocess(notebook, {'metadata': {'path': './'}})
    except Exception as e:
        # store traceback
        error = traceback.format_exc()
        print("Error during notebook execution:", e)

    # Save executed notebook
    notebook_content = nbformat.writes(notebook)

    # If we reach here, execution was successful
    return notebook_content, error



def append_result(a, b):
    """
    Appends two string, which might be produced by LLMs. E.g. in case the first thing contains ```python and the second
    starts with ```python, the beginning of the second string is removed.
    """
    if len(a) == 0:
        return b
    if len(b) == 0:
        return a

    possible_beginnings = POSSBILE_MARKDOWN_FENCES

    for beginning in possible_beginnings:
        if beginning in a and b.startswith(beginning + "\n"):
            b = b[len(beginning):]
            return a + b
    return a + b


def remove_ansi_escape_sequences(text):
    import re
    ansi_escape = re.compile(r'\x1B\[[0-?]*[ -/]*[@-~]')
    return ansi_escape.sub('', text)


def run_cli(command:str, check=False, verbose=False):
    import subprocess

    result = subprocess.run(command, shell=True, check=check, capture_output=True, text=True)
    if verbose:
        print("\n", result.stdout)
        print("\n", result.stderr)

    return f"## Command\n```\n{command}\n```\n## StdOut\n```\n{result.stdout}\n```\n## StdErr\n```\n{result.stderr}\n```\n"


def deploy(repository, issue, **kwargs):
    """
    Deploy the package to PyPI.

    Parameters
    ----------
    repository: str
        The repository name.
    issue: str
        The issue number where to post the deployment report.
    """
    #from ._github_utilities import add_comment_to_issue
    from ._ai_github_utilities import setup_ai_remark
    result1 = run_cli("python setup.py sdist bdist_wheel")
    result2 = run_cli("twine upload dist/*")
    Config.git_utilities.add_comment_to_issue(repository, issue, setup_ai_remark() + remove_ansi_escape_sequences(f"\n# Deployment report\n\n{result1}\n{result2}"))


def clean_output(repository, text):
    #from ._github_utilities import get_contributors
    # if all lines start with spaces (except 1st and last), remove those spaces in all lines
    lines = text.split("\n")
    while (True):
        if all([line.startswith(" ") for line in lines[1:]]):
            text = lines[0] + "\n" + "\n".join([line[1:] for line in lines[1:]])
            lines = text.split("\n")
        else:
            break

    # if text starts with ```markdown, remove that
    while text.startswith("\n") or text.startswith(" "):
        text = text.strip("\n").strip(" ")
    text = remove_outer_markdown(text)

    # if there are strangers tagged, remove those tags
    temp = text.split("```")
    for i in range(0, len(temp), 2):
        temp[i] = temp[i].replace("@", "@ ")
    text = "```".join(temp)
    contributors = Config.git_utilities.get_contributors(repository)
    for c in contributors:
        text = text.replace("@ " + c, "@" + c)
    return text


SENSIBLE_ENV_KEYS = ["ANTHROPIC_API_KEY",
                    "GOOGLE_API_KEY",
                    "OPENAI_API_KEY",
                    "GH_MODELS_API_KEY",
                    "MISTRAL_API_KEY",
                    "KISSKI_API_KEY",
                    "BLABLADOR_API_KEY",
                    "GITHUB_API_KEY",
                    "GITLAB_API_KEY",
                    "TWINE_USERNAME",
                    "TWINE_PASSWORD",
                    "HF_TOKEN",
                    "DEEPSEEK_API_KEY",
                    "CODECOV_TOKEN"]

def save_and_clear_environment():
    """Clear all environment variables and store the entire env in a dictionary for restoration later"""
    import os
    # Save the current environment
    saved_env = dict(os.environ)

    # Clear all environment variables
    for key in list(os.environ.keys()):
        if key in SENSIBLE_ENV_KEYS or \
            "password" in key.lower() or \
            "username" in key.lower() or \
            "key" in key.lower():
            del os.environ[key]

    return saved_env


def restore_environment(saved_env):
    """Restore an environment that was saved with save_and_clear_environment"""
    import os
    # Clear current environment
    for key in list(os.environ.keys()):
        del os.environ[key]

    # Restore saved environment
    os.environ.update(saved_env)


def redact_text(text):
    """Hide sensitive information from a string"""
    for key in SENSIBLE_ENV_KEYS:
        if key in os.environ and len(os.environ.get(key)) > 0:
            text = text.replace(os.environ.get(key), "***")
    return text


def file_list_from_commit_message_dict(repository, branch_name, commit_messages):
    list_of_links = []
    for k, v in commit_messages.items():

        if "https://github.com/" in Config.git_server_url:
            url_template = f"{Config.git_server_url}{repository}/blob/{branch_name}/"
        else:
            url_template = f"{Config.git_server_url}{repository}/-/blob/{branch_name}/"

        suffix = ""
        prefix = ""
        if any([k.endswith(f) for f in IMAGE_FILE_ENDINGS]):
            prefix = "!"
            if "https://github.com/" in Config.git_server_url:
                suffix = "?raw=true"
            else:
                url_template = url_template.replace("/blob/", "/raw/")

        list_of_links.append(f"{prefix}[{k}]({url_template}{k}{suffix})")
    return list_of_links


def ensure_images_shown(markdown, list_of_markdown_image_links):
    """
    If [bla](bla.png) in markdown, but not ![bla](bla.png), add the ! in front.
    """

    for i in list_of_markdown_image_links:
        if i[0] == "!":
            if i[1:] in markdown and i not in markdown:
                markdown = markdown.replace(i[1:], i)
    return markdown


def get_file_info(root_dir='.'):
    """Get a dictionary of files with their last change dates"""
    import os
    file_info = {}
    for root, dirs, files in os.walk(root_dir):
        for file in files:
            file_path = os.path.join(root, file)
            file_info[file_path] = os.path.getmtime(file_path)
    return file_info

def get_modified_files(old_file_info, root_dir='.'):
    """Get a list of modified files since get_file_info was called"""
    import os
    modified_files = []
    for root, dirs, files in os.walk(root_dir):
        for file in files:
            file_path = os.path.join(root, file)
            if file_path not in old_file_info:
                modified_files.append(file_path)  # New file
            elif os.path.getmtime(file_path) != old_file_info[file_path]:
                modified_files.append(file_path)  # Modified file

    result = []
    for m in modified_files:
        m = m.replace("\\", "/")
        if m.startswith("./"):
            m = m[2:]
        result.append(m)
    return result


def images_from_url_responses(response, input_shape=None):
    """Turns a list of OpenAI's URL responses into PIL images."""
    from skimage.io import imread
    from skimage import transform
    from PIL import Image
    # Removed unnecessary conversion to numpy arrays

    pil_images = [Image.fromarray(imread(item.url)) for item in response.data]

    if input_shape is not None:
        pil_images = [img.resize(input_shape) for img in pil_images]

    if len(pil_images) == 1:
        return pil_images[0]
    else:
        return pil_images


def make_slides(slides_description_json, filename="issue_slides.pptx"):
    """
    Create a PowerPoint slide deck from a JSON description of slides.

    Parameters
    ----------
    slides_description_json : str
        A JSON string representing the slides. Each slide is expected to be a dictionary with
        a 'title' key and a 'content' key, where 'content' is a list of strings.
    filename : str, optional
        The name of the file to save the presentation to. Defaults to 'issue_slides.pptx'.

    """
    import json
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pathlib import Path
    from PIL import Image
    import os
    import re

    print("SLIDES_JSON:", slides_description_json)

    # Parse json-encoded slide description
    slides_description_json = re.sub(r'[\x00-\x1f\x7f]', '', slides_description_json)

    slides_data = json.loads(slides_description_json)

    file_location = "slide_template.pptx"
    if not os.path.exists(file_location):
        file_location = Path(__file__).parent / "data" / "slide_template.pptx"

    # Create a presentation
    presentation = Presentation(file_location)

    # determine slide size
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Convert EMU to inches (1 inch = 914400 EMUs)
    width_inch = slide_width / 914400
    height_inch = slide_height / 914400

    top = Inches(2)
    bottom = Inches(1)

    # Iterate through slide data to create slides
    for i, slide_data in enumerate(slides_data):
        # Add a slide with title and content layout
        slide_layout = presentation.slide_layouts[0 if i == 0 else 1] # choose first or second layout
        slide = presentation.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.title
        title_box.text = slide_data['title']

        # Calculate width for content columns
        num_columns = len(slide_data['content'])

        # remove all placeholders except the title
        for shape in slide.placeholders:
            if shape != title_box:
                #shape.text = ""
                top = shape.top
                bottom = height_inch - top
                slide.shapes._spTree.remove(shape._element)

        for i, content in enumerate(slide_data['content']):
            content_width = (Inches(width_inch) - top) / num_columns
            content_height = (Inches(height_inch) - top) - bottom
            left = Inches(1) + i * (content_width + Inches(0.1))

            print("Left", left)

            if any([content.endswith(f) for f in IMAGE_FILE_ENDINGS]):
                image_path = content

                # load image using PIL and determine width/height ratio
                image = Image.open(image_path)
                width, height = image.size
                aspect_ratio = width / height

                if content_width / aspect_ratio > content_height - bottom:
                    content_width = (content_height - bottom) / aspect_ratio

                content_box = slide.shapes.add_picture(image_path, left=left, top=top, width=content_width, height=content_width / aspect_ratio)
            else:
                content_box = slide.shapes.add_textbox(left=left, top=top, width=content_width, height=content_height)
                text_frame = content_box.text_frame
                text_frame.text = content
                text_frame.word_wrap = True

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Set the font size for each run in each paragraph
                        run.font.size = Pt(24)  # Set font size to 24 points


    # Save presentation
    presentation.save(filename)


def read_text_file(filename):
    """Read the content of a text file."""
    with open(filename, 'r') as file:
        return file.read()

def write_text_file(filename, content):
    """Write content to a text file."""
    with open(filename, 'w') as file:
        file.write(content)

def read_binary_file(filename):
    """Read the content of a binary file."""
    with open(filename, 'rb') as file:
        return file.read()

def write_binary_file(filename, content):
    """Write content to a binary file."""
    with open(filename, 'wb') as file:
        file.write(content)
